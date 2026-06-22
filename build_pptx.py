#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_pptx.py — Linuxコマンド道場のレッスン内容を PowerPoint(.pptx) に変換する
ターミナル風のデザインで、各コマンドの説明・使用例・ミッションを1枚ずつスライド化する。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- カラーパレット（Webサイトと統一） ----------
BG_DARK     = RGBColor(0x0F, 0x17, 0x2A)  # ターミナル背景
PANEL       = RGBColor(0x13, 0x1C, 0x31)
TEXT_LIGHT  = RGBColor(0xE2, 0xE8, 0xF0)
MUTED       = RGBColor(0x94, 0xA3, 0xB8)
ACCENT      = RGBColor(0x81, 0x8C, 0xF8)  # インディゴ
GREEN       = RGBColor(0x4A, 0xDE, 0x80)
BLUE        = RGBColor(0x60, 0xA5, 0xFA)
RED         = RGBColor(0xF8, 0x71, 0x71)
YELLOW      = RGBColor(0xFB, 0xBD, 0x2E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

FONT_BODY = "Meiryo"          # 日本語が崩れにくいフォント
FONT_MONO = "Consolas"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------- レッスンデータ（lessons.js と対応） ----------
LESSONS = [
    {
        "icon": "📂", "no": "01", "title": "ファイルの一覧を見る",
        "cmd": "ls",
        "desc": "今いる場所にあるファイルやディレクトリ（フォルダ）の一覧を表示する、最もよく使うコマンドです。",
        "examples": [
            ("ls", "一覧を表示"),
            ("ls -l", "詳細（権限・サイズ）付きで表示"),
            ("ls -a", "隠しファイルも表示"),
        ],
        "demo": ["user@dojo:~$ ls",
                 "documents/  projects/  memo.txt  welcome.txt"],
        "mission": "ターミナルで ls を実行して、ホームにあるファイルを表示してみよう。",
    },
    {
        "icon": "🧭", "no": "02", "title": "今どこ？ 移動する",
        "cmd": "pwd / cd",
        "desc": "pwd で現在地の確認、cd でディレクトリ間を移動します。Linux はフォルダが階層構造になっています。",
        "examples": [
            ("pwd", "現在地の絶対パスを表示"),
            ("cd documents", "documents の中へ移動"),
            ("cd ..", "一つ上の階層へ戻る"),
            ("cd ~", "ホームディレクトリへ戻る"),
        ],
        "demo": ["user@dojo:~$ pwd",
                 "/home/user",
                 "user@dojo:~$ cd documents",
                 "user@dojo:~/documents$"],
        "mission": "cd documents で documents ディレクトリへ移動しよう。",
    },
    {
        "icon": "📄", "no": "03", "title": "ファイルの中身を見る",
        "cmd": "cat",
        "desc": "cat はファイルの中身を画面に表示します。head / tail で先頭・末尾だけ見ることもできます。",
        "examples": [
            ("cat welcome.txt", "ファイルの中身を表示"),
            ("head todo.md", "先頭10行だけ表示"),
            ("tail todo.md", "末尾10行だけ表示"),
        ],
        "demo": ["user@dojo:~$ cat memo.txt",
                 "買い物リスト",
                 "- りんご",
                 "- ぱん",
                 "- ぎゅうにゅう"],
        "mission": "cat memo.txt で買い物メモを表示しよう。",
    },
    {
        "icon": "🛠️", "no": "04", "title": "ディレクトリとファイルを作る",
        "cmd": "mkdir / touch",
        "desc": "mkdir で新しいフォルダ、touch で空のファイルを作成します。",
        "examples": [
            ("mkdir practice", "practice フォルダを作成"),
            ("touch note.txt", "空の note.txt を作成"),
            ("echo \"やあ\" > hi.txt", "文字を書き込みつつ作成"),
        ],
        "demo": ["user@dojo:~$ mkdir practice",
                 "user@dojo:~$ ls",
                 "documents/  practice/  memo.txt  welcome.txt"],
        "mission": "mkdir practice で practice という新しいフォルダを作ろう。",
    },
    {
        "icon": "✍️", "no": "05", "title": "ファイルに書き込む",
        "cmd": "echo >",
        "desc": "echo は文字を表示。> （リダイレクト）と組み合わせるとファイルに保存、>> なら末尾に追記します。",
        "examples": [
            ("echo こんにちは", "画面に表示するだけ"),
            ("echo こんにちは > greet.txt", "greet.txt に書き込み"),
            ("echo またね >> greet.txt", "greet.txt に追記"),
        ],
        "demo": ["user@dojo:~$ echo hello > greet.txt",
                 "user@dojo:~$ cat greet.txt",
                 "hello"],
        "mission": "echo hello > greet.txt で greet.txt を作って書き込もう。",
    },
    {
        "icon": "🔁", "no": "06", "title": "コピー・移動・削除",
        "cmd": "cp / mv / rm",
        "desc": "ファイルのコピー・移動（リネーム）・削除を行います。rm したファイルは元に戻らないので注意！",
        "examples": [
            ("cp a.txt b.txt", "a.txt を b.txt にコピー"),
            ("mv old.txt new.txt", "名前を変更（リネーム）"),
            ("rm a.txt", "ファイルを削除"),
            ("rm -r folder", "フォルダごと削除"),
        ],
        "demo": ["user@dojo:~$ cp welcome.txt hello.txt",
                 "user@dojo:~$ ls",
                 "hello.txt  memo.txt  welcome.txt"],
        "mission": "cp welcome.txt hello.txt で welcome.txt をコピーしよう。",
    },
    {
        "icon": "🔍", "no": "07", "title": "検索と集計",
        "cmd": "grep / wc",
        "desc": "grep はファイル内から特定の文字を含む行を検索、wc は行数や単語数を数えます。",
        "examples": [
            ("grep りんご memo.txt", "「りんご」を含む行を表示"),
            ("wc -l memo.txt", "行数を数える"),
            ("tree", "フォルダ構造をツリー表示"),
        ],
        "demo": ["user@dojo:~$ grep ぱん memo.txt",
                 "- ぱん",
                 "user@dojo:~$ wc -l memo.txt",
                 "4 memo.txt"],
        "mission": "grep ぱん memo.txt で「ぱん」を含む行を探そう。",
    },
    {
        "icon": "🎓", "no": "08", "title": "総合演習 — 卒業ミッション",
        "cmd": "総合演習",
        "desc": "学んだコマンドを組み合わせて、バックアップ作業に挑戦しましょう。これができれば卒業です！",
        "examples": [
            ("cd ~", "ホームへ移動"),
            ("mkdir backup", "backup フォルダを作成"),
            ("cp documents/report.txt backup/", "レポートをコピー"),
            ("ls backup", "結果を確認"),
        ],
        "demo": ["user@dojo:~$ mkdir backup",
                 "user@dojo:~$ cp documents/report.txt backup/",
                 "user@dojo:~$ ls backup",
                 "report.txt"],
        "mission": "backup フォルダを作り、その中に report.txt をコピーしよう。",
    },
]

# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, runs, size=18, bold=False, color=TEXT_LIGHT,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0):
    """runs: 文字列 もしくは [(text, color, bold, font), ...] のリスト"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, color, bold, font)]
    for item in runs:
        text, c, b, f = (item + (None,) * 4)[:4]
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = b if b is not None else bold
        r.font.color.rgb = c if c is not None else color
        r.font.name = f if f else font
    return tb


def fill_background(slide, color):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)


# ---------- 表紙 ----------
def title_slide():
    s = prs.slides.add_slide(BLANK)
    fill_background(s, BG_DARK)
    # アクセントの帯
    add_rect(s, 0, Inches(3.05), SLIDE_W, Inches(0.06), ACCENT)
    add_text(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
             "🐧  Linuxコマンド道場", size=54, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.25), Inches(11.3), Inches(0.8),
             "ターミナルで実際に打ちながら学ぶ 基本コマンド入門", size=22,
             color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.6),
             "全8レッスン / ls・cd・cat・mkdir・cp・grep ほか", size=16,
             color=GREEN, font=FONT_MONO, align=PP_ALIGN.CENTER)


# ---------- 目次 ----------
def toc_slide():
    s = prs.slides.add_slide(BLANK)
    fill_background(s, BG_DARK)
    add_rect(s, Inches(0.7), Inches(0.5), Inches(0.12), Inches(0.7), ACCENT)
    add_text(s, Inches(1.0), Inches(0.5), Inches(11), Inches(0.8),
             "目次 — Contents", size=32, bold=True, color=WHITE)
    col_x = [Inches(1.0), Inches(7.0)]
    for i, L in enumerate(LESSONS):
        col = 0 if i < 4 else 1
        row = i % 4
        y = Inches(1.7 + row * 1.25)
        card = add_rect(s, col_x[col], y, Inches(5.5), Inches(1.0), PANEL)
        add_text(s, col_x[col] + Inches(0.2), y + Inches(0.12), Inches(0.9), Inches(0.8),
                 f"{L['icon']}", size=28, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, col_x[col] + Inches(1.1), y + Inches(0.14), Inches(4.2), Inches(0.4),
                 [(f"LESSON {L['no']}  ", ACCENT, True, FONT_MONO),
                  (L['cmd'], GREEN, True, FONT_MONO)], size=13)
        add_text(s, col_x[col] + Inches(1.1), y + Inches(0.5), Inches(4.2), Inches(0.4),
                 L['title'], size=15, bold=True, color=TEXT_LIGHT)


# ---------- 各レッスン ----------
def lesson_slide(L):
    s = prs.slides.add_slide(BLANK)
    fill_background(s, BG_DARK)

    # ヘッダー
    add_text(s, Inches(0.7), Inches(0.45), Inches(8), Inches(0.5),
             [(f"LESSON {L['no']}", ACCENT, True, FONT_MONO)], size=16)
    add_text(s, Inches(0.7), Inches(0.85), Inches(9), Inches(0.8),
             f"{L['icon']}  {L['title']}", size=30, bold=True, color=WHITE)
    # コマンドバッジ
    badge = add_rect(s, Inches(10.0), Inches(0.6), Inches(2.6), Inches(0.75), PANEL,
                     line=ACCENT)
    add_text(s, Inches(10.0), Inches(0.6), Inches(2.6), Inches(0.75),
             L['cmd'], size=22, bold=True, color=GREEN, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 説明文
    add_text(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(0.9),
             L['desc'], size=16, color=TEXT_LIGHT, line_spacing=1.25)

    # 左：使用例テーブル
    lx, ly = Inches(0.7), Inches(2.85)
    add_text(s, lx, ly - Inches(0.05), Inches(5.5), Inches(0.4),
             "よく使う形", size=15, bold=True, color=YELLOW)
    for i, (cmd, desc) in enumerate(L['examples']):
        ry = ly + Inches(0.5 + i * 0.78)
        add_rect(s, lx, ry, Inches(5.6), Inches(0.66), PANEL)
        add_text(s, lx + Inches(0.15), ry + Inches(0.06), Inches(5.3), Inches(0.32),
                 cmd, size=14, bold=True, color=GREEN, font=FONT_MONO)
        add_text(s, lx + Inches(0.15), ry + Inches(0.36), Inches(5.3), Inches(0.28),
                 desc, size=11.5, color=MUTED)

    # 右：ターミナルデモ
    tx, ty = Inches(6.7), Inches(2.85)
    add_text(s, tx, ty - Inches(0.05), Inches(5.9), Inches(0.4),
             "実行イメージ", size=15, bold=True, color=YELLOW)
    term = add_rect(s, tx, ty + Inches(0.45), Inches(5.95), Inches(2.7),
                    RGBColor(0x0B, 0x11, 0x20), line=RGBColor(0x1E, 0x29, 0x3B))
    # タイトルバーのドット
    for i, c in enumerate([RED, YELLOW, GREEN]):
        add_rect(s, tx + Inches(0.2 + i * 0.28), ty + Inches(0.62),
                 Inches(0.16), Inches(0.16), c)
    # デモ本文
    tb = s.shapes.add_textbox(tx + Inches(0.2), ty + Inches(0.95),
                              Inches(5.6), Inches(2.05))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(L['demo']):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        is_cmd = line.startswith("user@dojo")
        if is_cmd and "$ " in line:
            prompt, _, rest = line.partition("$ ")
            r1 = p.add_run(); r1.text = prompt + "$ "
            r1.font.color.rgb = GREEN; r1.font.size = Pt(13); r1.font.name = FONT_MONO
            r2 = p.add_run(); r2.text = rest
            r2.font.color.rgb = TEXT_LIGHT; r2.font.size = Pt(13); r2.font.name = FONT_MONO
        else:
            r = p.add_run(); r.text = line
            r.font.color.rgb = TEXT_LIGHT; r.font.size = Pt(13); r.font.name = FONT_MONO

    # 下：ミッション帯
    my = Inches(5.95)
    add_rect(s, Inches(0.7), my, Inches(11.93), Inches(0.95),
             RGBColor(0x1E, 0x27, 0x45), line=ACCENT)
    add_text(s, Inches(0.95), my + Inches(0.1), Inches(11.4), Inches(0.35),
             [("📝 ミッション", YELLOW, True, FONT_BODY)], size=15)
    add_text(s, Inches(0.95), my + Inches(0.46), Inches(11.4), Inches(0.4),
             L['mission'], size=14, color=TEXT_LIGHT)


# ---------- まとめ ----------
def summary_slide():
    s = prs.slides.add_slide(BLANK)
    fill_background(s, BG_DARK)
    add_text(s, Inches(0.7), Inches(0.5), Inches(11), Inches(0.8),
             "🎉 おさらい — チートシート", size=32, bold=True, color=WHITE)
    cheats = [
        ("ls", "一覧表示"), ("pwd", "現在地"), ("cd", "移動"),
        ("cat", "中身表示"), ("echo >", "書き込み"), ("mkdir", "フォルダ作成"),
        ("touch", "ファイル作成"), ("cp", "コピー"), ("mv", "移動/改名"),
        ("rm", "削除"), ("grep", "検索"), ("wc", "行数集計"),
    ]
    for i, (cmd, desc) in enumerate(cheats):
        col = i % 3
        row = i // 3
        x = Inches(0.7 + col * 4.1)
        y = Inches(1.7 + row * 1.15)
        add_rect(s, x, y, Inches(3.85), Inches(0.95), PANEL)
        add_text(s, x + Inches(0.2), y + Inches(0.13), Inches(3.5), Inches(0.4),
                 cmd, size=18, bold=True, color=GREEN, font=FONT_MONO)
        add_text(s, x + Inches(0.2), y + Inches(0.53), Inches(3.5), Inches(0.34),
                 desc, size=13, color=MUTED)
    add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.6),
             "ブラウザ版『Linuxコマンド道場』で、実際に打って練習しよう！ 🐧",
             size=16, color=ACCENT, align=PP_ALIGN.CENTER)


# ---------- 生成 ----------
title_slide()
toc_slide()
for L in LESSONS:
    lesson_slide(L)
summary_slide()

OUT = "Linuxコマンド道場.pptx"
prs.save(OUT)
print(f"作成完了: {OUT}  （全{len(prs.slides._sldIdLst)}スライド）")
